from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ────────────────────────────────────
NAVY     = RGBColor(0x0D, 0x2B, 0x55)
BLUE     = RGBColor(0x16, 0x5A, 0xB5)
LBLUE    = RGBColor(0x3B, 0x82, 0xF6)
ICE      = RGBColor(0xEB, 0xF2, 0xFF)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x1E, 0x29, 0x3B)
GRAY     = RGBColor(0x47, 0x55, 0x6B)
PALE     = RGBColor(0xB0, 0xC8, 0xFF)

SHOTS = r"c:\Users\geori\OneDrive\Documentos\projects\my-first-project\screenshots"
OUT   = r"c:\Users\geori\OneDrive\Documentos\projects\my-first-project\Vanguard_Clothier_ERP_Presentation.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── Helpers ────────────────────────────────────

def rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background() if not line else None
    if line: s.line.color.rgb = line
    return s

def add_text(slide, x, y, w, h, text, size, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return box

def add_para(tf, text, size, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, before=0):
    p = tf.add_paragraph(); p.alignment = align; p.space_before = Pt(before)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return p

def add_img(slide, name, x, y, w, h):
    path = os.path.join(SHOTS, name)
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)

def base(slide):
    rect(slide, 0, 0, W, H, WHITE)
    rect(slide, 0, 0, W, Inches(0.5), NAVY)
    rect(slide, 0, H - Inches(0.28), W, Inches(0.28), NAVY)

def header(slide, title):
    base(slide)
    rect(slide, 0, Inches(0.5), W, Inches(0.85), ICE)
    rect(slide, 0, Inches(0.5), Inches(0.1), Inches(0.85), BLUE)
    add_text(slide, Inches(0.25), Inches(0.56), W - Inches(0.4), Inches(0.75),
             title, 25, bold=True, color=NAVY)
    rect(slide, 0, Inches(1.35), W, Inches(0.035), BLUE)
    return Inches(1.45)   # content Y start

def caption(slide, text, x, y, w):
    rect(slide, x, y, w, Inches(0.32), NAVY)
    add_text(slide, x + Inches(0.1), y + Inches(0.03), w - Inches(0.15),
             Inches(0.28), text, 12, bold=True, color=WHITE)


# ══════════════════════════════════════════════
# СЛАЙД 1 — Титульный
# ══════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])
base(s1)
rect(s1, 0, 0, W, Inches(2.1), NAVY)
rect(s1, 0, Inches(2.1), W, Inches(0.07), BLUE)
rect(s1, 0, Inches(2.17), W, H - Inches(2.45), ICE)

b1 = s1.shapes.add_textbox(Inches(0.5), Inches(0.15), W - Inches(1.0), Inches(1.0))
tf1 = b1.text_frame; tf1.word_wrap = True
p = tf1.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = ("ФГБОУ ВО «Северо-Кавказский горно-металлургический институт (ГТУ)»\n"
          "Многопрофильный профессиональный колледж СКГМИ (ГТУ)")
r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = PALE

b2 = s1.shapes.add_textbox(Inches(0.5), Inches(1.0), W - Inches(1.0), Inches(1.15))
tf2 = b2.text_frame; tf2.word_wrap = True
p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = ("Разработка автоматизированной информационной системы\n"
           "учета товаров в магазине одежды")
r2.font.size = Pt(29); r2.font.bold = True; r2.font.color.rgb = WHITE

# Логотип V
rect(s1, W/2 - Inches(0.48), Inches(2.3), Inches(0.96), Inches(0.96), BLUE)
add_text(s1, W/2 - Inches(0.48), Inches(2.32), Inches(0.96), Inches(0.9),
         "V", 38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

b3 = s1.shapes.add_textbox(Inches(3.5), Inches(3.45), Inches(6.5), Inches(2.0))
tf3 = b3.text_frame; tf3.word_wrap = True
for i, (lbl, val) in enumerate([
    ("Руководитель:", "Цогоева З.О., преподаватель высшей категории"),
    ("Студент:",      "Созанов Георгий Валерьевич"),
    ("Группа:",       "ИС-22-2"),
]):
    p3 = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(8)
    ra = p3.add_run(); ra.text = f"{lbl}  "
    ra.font.size = Pt(17); ra.font.bold = True; ra.font.color.rgb = LBLUE
    rb = p3.add_run(); rb.text = val
    rb.font.size = Pt(17); rb.font.bold = True; rb.font.color.rgb = DARK

add_text(s1, 0, H - Inches(0.62), W, Inches(0.38),
         "Владикавказ, 2026", 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# СЛАЙД 2 — Актуальность + Цель (объединён)
# ══════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
cy = header(s2, "Актуальность и цель работы")

# Левая колонка — Проблема
lw = Inches(6.3)
rect(s2, Inches(0.3), cy, lw, Inches(0.44), NAVY)
add_text(s2, Inches(0.4), cy + Inches(0.05), lw, Inches(0.38),
         "Проблема", 15, bold=True, color=WHITE)

lb = s2.shapes.add_textbox(Inches(0.3), cy + Inches(0.5), lw, Inches(4.5))
ltf = lb.text_frame; ltf.word_wrap = True
for b in [
    "Ручной учёт товаров: ошибки, пересортица, финансовые потери",
    "Нет контроля складских остатков в реальном времени",
    "Отсутствие инструментов анализа продаж и отчётности",
    "Нет онлайн-канала продаж для привлечения клиентов",
]:
    p = ltf.paragraphs[0] if not ltf.paragraphs[0].runs else ltf.add_paragraph()
    p.space_before = Pt(10)
    ra = p.add_run(); ra.text = "▸  "; ra.font.size = Pt(16)
    ra.font.bold = True; ra.font.color.rgb = LBLUE
    rb = p.add_run(); rb.text = b
    rb.font.size = Pt(16); rb.font.color.rgb = DARK

# Разделитель
rect(s2, Inches(6.75), cy, Inches(0.04), Inches(5.1), BLUE)

# Правая колонка — Цель
rw = Inches(6.1); rx = Inches(6.95)
rect(s2, rx, cy, rw, Inches(0.44), BLUE)
add_text(s2, rx + Inches(0.1), cy + Inches(0.05), rw, Inches(0.38),
         "Цель работы", 15, bold=True, color=WHITE)

rb2 = s2.shapes.add_textbox(rx, cy + Inches(0.5), rw, Inches(1.4))
rtf = rb2.text_frame; rtf.word_wrap = True
p = rtf.paragraphs[0]
r = p.add_run()
r.text = ("Разработать АИС для автоматизации учёта товаров "
          "и управления бизнес-процессами магазина одежды")
r.font.size = Pt(16); r.font.color.rgb = DARK

rect(s2, rx, cy + Inches(1.95), rw, Inches(0.44), NAVY)
add_text(s2, rx + Inches(0.1), cy + Inches(2.0), rw, Inches(0.38),
         "Объект и предмет", 15, bold=True, color=WHITE)

rb3 = s2.shapes.add_textbox(rx, cy + Inches(2.45), rw, Inches(1.8))
rtf2 = rb3.text_frame; rtf2.word_wrap = True
for obj, val in [
    ("Объект:", "процессы учёта и управления в магазине одежды"),
    ("Предмет:", "программные средства автоматизации учёта товаров"),
]:
    p = rtf2.paragraphs[0] if not rtf2.paragraphs[0].runs else rtf2.add_paragraph()
    p.space_before = Pt(10)
    ra = p.add_run(); ra.text = f"{obj}  "; ra.font.size = Pt(15)
    ra.font.bold = True; ra.font.color.rgb = BLUE
    rb = p.add_run(); rb.text = val
    rb.font.size = Pt(15); rb.font.color.rgb = DARK


# ══════════════════════════════════════════════
# СЛАЙД 3 — Задачи
# ══════════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6])
cy = header(s3, "Задачи работы")

tasks = [
    ("1", "Анализ предметной области",
          "Изучить бизнес-процессы магазина одежды, определить требования к системе"),
    ("2", "Проектирование БД",
          "Разработать схему данных: 15+ моделей с каскадными связями (Prisma ORM)"),
    ("3", "Модуль управления складом",
          "Реализовать полный CRUD товаров, вариаций, категорий и брендов"),
    ("4", "POS-касса и продажи",
          "Оформление розничных продаж, скидки, чек со штрихкодом, Z-отчёт"),
    ("5", "CRM и программа лояльности",
          "Клиентская база, история покупок, начисление и списание баллов"),
    ("6", "Финансы и аналитика",
          "P&L отчёт, расходы, топ товаров, экспорт в Excel"),
    ("7", "Онлайн-витрина",
          "Публичный каталог, корзина, оформление заказа, личный кабинет"),
]

col_gap = Inches(0.25)
cw = (W - Inches(0.6) - col_gap) / 2
sy = cy + Inches(0.1)
row_h = Inches(0.75)

for i, (num, title, desc) in enumerate(tasks):
    col = i % 2
    row = i // 2
    x = Inches(0.3) + col * (cw + col_gap)
    y = sy + row * (row_h + Inches(0.12))
    if y + row_h > H - Inches(0.35):
        break
    rect(s3, x, y, cw, row_h, ICE, LBLUE)
    rect(s3, x, y, Inches(0.38), row_h, BLUE)
    add_text(s3, x + Inches(0.05), y + Inches(0.15), Inches(0.3), Inches(0.45),
             num, 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb = s3.shapes.add_textbox(x + Inches(0.45), y + Inches(0.06),
                                cw - Inches(0.52), row_h - Inches(0.1))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = NAVY
    add_para(tf, desc, 12, color=GRAY, before=3)


# ══════════════════════════════════════════════
# СЛАЙД 4 — Архитектура
# ══════════════════════════════════════════════
s4 = prs.slides.add_slide(prs.slide_layouts[6])
cy = header(s4, "Архитектура системы")

cols = [
    ("Фронтенд",    ["React 19", "TypeScript", "Tailwind CSS", "Framer Motion"]),
    ("Бэкенд",      ["Node.js", "Express.js", "REST API", "JWT Auth"]),
    ("База данных", ["SQLite", "Prisma ORM", "15+ моделей", "PM2 деплой"]),
]
cw = Inches(3.6); gap = Inches(0.27); ch = Inches(3.7)
sx = Inches(0.38); sy = cy + Inches(0.15)

for i, (title, items) in enumerate(cols):
    x = sx + i * (cw + gap)
    rect(s4, x, sy, cw, ch, ICE, BLUE)
    rect(s4, x, sy, cw, Inches(0.54), NAVY)
    add_text(s4, x, sy + Inches(0.07), cw, Inches(0.46),
             title, 17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    ib = s4.shapes.add_textbox(x + Inches(0.15), sy + Inches(0.65),
                                cw - Inches(0.3), ch - Inches(0.75))
    itf = ib.text_frame; itf.word_wrap = True
    for j, item in enumerate(items):
        ip = itf.paragraphs[0] if j == 0 else itf.add_paragraph()
        ip.alignment = PP_ALIGN.CENTER; ip.space_before = Pt(14)
        ir = ip.add_run(); ir.text = item
        ir.font.size = Pt(16); ir.font.color.rgb = DARK
    if i < 2:
        add_text(s4, x + cw + Inches(0.01), sy + ch/2 - Inches(0.28),
                 gap + Inches(0.06), Inches(0.5),
                 "⟺", 22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

rect(s4, Inches(0.38), sy + ch + Inches(0.2), W - Inches(0.7), Inches(0.5), NAVY)
add_text(s4, Inches(0.38), sy + ch + Inches(0.21), W - Inches(0.7), Inches(0.5),
         "Роли пользователей:    ADMIN   /   SELLER   /   STOREKEEPER",
         17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# СЛАЙД 5 — Модули
# ══════════════════════════════════════════════
s5 = prs.slides.add_slide(prs.slide_layouts[6])
cy = header(s5, "Функциональные модули системы")

modules = [
    ("Склад",       "Учёт товаров, вариации (размер/цвет/цена), категории, бренды"),
    ("POS-касса",   "Оформление продаж, скидки, чек со штрихкодом, Z-отчёт"),
    ("Заказы",      "Онлайн-заказы, управление статусами, обработка возвратов"),
    ("CRM",         "База клиентов, программа лояльности, история покупок"),
    ("Финансы",     "Доходы/расходы, P&L с фильтром по периоду"),
    ("Аналитика",   "Топ товаров, low-stock, журнал операций, экспорт Excel"),
    ("Витрина",     "Публичный каталог, корзина, оформление заказа онлайн"),
]
rh = Inches(0.57); c1w = Inches(2.8); c2w = Inches(9.85)
sx = Inches(0.35); sy = cy + Inches(0.05)

rect(s5, sx, sy, c1w, rh, NAVY)
rect(s5, sx + c1w, sy, c2w, rh, BLUE)
add_text(s5, sx + Inches(0.1), sy + Inches(0.1), c1w, rh,
         "Модуль", 14, bold=True, color=WHITE)
add_text(s5, sx + c1w + Inches(0.1), sy + Inches(0.1), c2w, rh,
         "Функциональность", 14, bold=True, color=WHITE)

for i, (mod, desc) in enumerate(modules):
    y = sy + rh * (i + 1)
    bg = ICE if i % 2 == 0 else WHITE
    rect(s5, sx,       y, c1w, rh, bg, PALE)
    rect(s5, sx + c1w, y, c2w, rh, bg, PALE)
    add_text(s5, sx + Inches(0.1), y + Inches(0.1), c1w, rh,
             mod, 13, bold=True, color=BLUE)
    add_text(s5, sx + c1w + Inches(0.1), y + Inches(0.1), c2w, rh,
             desc, 13, color=DARK)


# ══════════════════════════════════════════════
# СЛАЙД 6 — Демо: Dashboard + POS
# ══════════════════════════════════════════════
s6 = prs.slides.add_slide(prs.slide_layouts[6])
cy = header(s6, "Демонстрация: Dashboard и POS-касса")

iw = Inches(6.2); ih = Inches(2.6)
lx = Inches(0.25); rx = Inches(6.65)
iy = cy + Inches(0.1)

add_img(s6, "dashboard.png", lx, iy,        iw, ih)
add_img(s6, "pos.png",       lx, iy + ih + Inches(0.15), iw, ih)
caption(s6, "Dashboard — сводная панель системы",  lx, iy + ih - Inches(0.01), iw)
caption(s6, "POS-касса — оформление розничной продажи", lx, iy + ih*2 + Inches(0.14), iw)

db = s6.shapes.add_textbox(rx, iy, W - rx - Inches(0.15), H - iy - Inches(0.35))
dtf = db.text_frame; dtf.word_wrap = True
for title, desc in [
    ("Dashboard",   "Выручка, остатки, ожидающие заказы, журнал операций в реальном времени"),
    ("POS-касса",   "Поиск товара по артикулу, фильтр по категориям, поле скидки, ограничение по остатку"),
    ("Чек",         "Автоматическая печать чека со штрихкодом через react-barcode, реквизиты магазина"),
    ("CommandPalette", "Быстрая навигация Ctrl+K с фильтром команд по роли пользователя"),
]:
    p = dtf.paragraphs[0] if not dtf.paragraphs[0].runs else dtf.add_paragraph()
    p.space_before = Pt(10)
    ra = p.add_run(); ra.text = f"{title}:  "; ra.font.size = Pt(14)
    ra.font.bold = True; ra.font.color.rgb = NAVY
    rb = p.add_run(); rb.text = desc
    rb.font.size = Pt(13); rb.font.color.rgb = DARK


# ══════════════════════════════════════════════
# СЛАЙД 7 — Демо: Склад + CRM
# ══════════════════════════════════════════════
s7 = prs.slides.add_slide(prs.slide_layouts[6])
cy = header(s7, "Демонстрация: Склад и CRM")

iy = cy + Inches(0.1)
add_img(s7, "inventory.png", lx, iy,        iw, ih)
add_img(s7, "crm.png",       lx, iy + ih + Inches(0.15), iw, ih)
caption(s7, "Склад — управление товарным ассортиментом",   lx, iy + ih - Inches(0.01), iw)
caption(s7, "CRM — клиентская база и программа лояльности", lx, iy + ih*2 + Inches(0.14), iw)

db = s7.shapes.add_textbox(rx, iy, W - rx - Inches(0.15), H - iy - Inches(0.35))
dtf = db.text_frame; dtf.word_wrap = True
for title, desc in [
    ("Товары",       "Полный CRUD: создание, редактирование, удаление товаров и категорий"),
    ("Вариации",     "Привязка размеров, цветов и цен к каждой позиции, контроль остатков"),
    ("Клиенты",      "Регистрация клиентов, поиск по ФИО / телефону / email"),
    ("Лояльность",   "Начисление баллов при покупке, списание при возврате, история транзакций"),
    ("Фильтрация",   "Поиск и сортировка по бренду, категории, остатку, цене"),
]:
    p = dtf.paragraphs[0] if not dtf.paragraphs[0].runs else dtf.add_paragraph()
    p.space_before = Pt(10)
    ra = p.add_run(); ra.text = f"{title}:  "; ra.font.size = Pt(14)
    ra.font.bold = True; ra.font.color.rgb = NAVY
    rb = p.add_run(); rb.text = desc
    rb.font.size = Pt(13); rb.font.color.rgb = DARK


# ══════════════════════════════════════════════
# СЛАЙД 8 — Демо: Витрина + Заказы
# ══════════════════════════════════════════════
s8 = prs.slides.add_slide(prs.slide_layouts[6])
cy = header(s8, "Демонстрация: Онлайн-витрина и Заказы")

iy = cy + Inches(0.1)
add_img(s8, "storefront.png", lx, iy,        iw, ih)
add_img(s8, "orders.png",     lx, iy + ih + Inches(0.15), iw, ih)
caption(s8, "Онлайн-витрина — публичный каталог товаров", lx, iy + ih - Inches(0.01), iw)
caption(s8, "Заказы — управление статусами",               lx, iy + ih*2 + Inches(0.14), iw)

db = s8.shapes.add_textbox(rx, iy, W - rx - Inches(0.15), H - iy - Inches(0.35))
dtf = db.text_frame; dtf.word_wrap = True
for title, desc in [
    ("Каталог",     "Фильтрация по цене, категориям и скидкам; карточки с фото и вариациями"),
    ("Корзина",     "Добавление товаров → оформление заказа → автосписание остатка склада"),
    ("Статусы",     "PENDING → CONFIRMED → SHIPPED → DELIVERED с уведомлениями"),
    ("Возврат",     "Обработка возвратов с автовосстановлением складского остатка"),
    ("Кабинет",     "История заказов, баллы лояльности, информационные страницы"),
]:
    p = dtf.paragraphs[0] if not dtf.paragraphs[0].runs else dtf.add_paragraph()
    p.space_before = Pt(10)
    ra = p.add_run(); ra.text = f"{title}:  "; ra.font.size = Pt(14)
    ra.font.bold = True; ra.font.color.rgb = NAVY
    rb = p.add_run(); rb.text = desc
    rb.font.size = Pt(13); rb.font.color.rgb = DARK


# ══════════════════════════════════════════════
# СЛАЙД 9 — Тестирование + БД + Результаты
# ══════════════════════════════════════════════
s9 = prs.slides.add_slide(prs.slide_layouts[6])
cy = header(s9, "Тестирование, база данных и результаты")

# Left: Testing + DB
lw = Inches(6.0)
add_text(s9, Inches(0.3), cy + Inches(0.05), lw, Inches(0.38),
         "Тестирование системы", 15, bold=True, color=NAVY)

lb = s9.shapes.add_textbox(Inches(0.3), cy + Inches(0.48), lw, Inches(2.2))
ltf = lb.text_frame; ltf.word_wrap = True
for e in [
    "Unit-тесты (Jest): 28 тестов — расчёт чека, скидки, баллы лояльности, возвраты, остатки",
    "Ручное тестирование: 15+ сценариев (POS, возврат, отмена заказа, фильтры)",
    "Тест безопасности: попытка подмены цены с клиента — сервер отклоняет",
    "Тест транзакции: одновременное списание остатка и начисление баллов — атомарно",
]:
    p = ltf.paragraphs[0] if not ltf.paragraphs[0].runs else ltf.add_paragraph()
    p.space_before = Pt(7)
    ra = p.add_run(); ra.text = "✓  "; ra.font.size = Pt(13)
    ra.font.bold = True; ra.font.color.rgb = LBLUE
    rb = p.add_run(); rb.text = e
    rb.font.size = Pt(13); rb.font.color.rgb = DARK

add_text(s9, Inches(0.3), cy + Inches(2.75), lw, Inches(0.38),
         "Структура базы данных (15+ моделей)", 15, bold=True, color=NAVY)

lb2 = s9.shapes.add_textbox(Inches(0.3), cy + Inches(3.18), lw, Inches(2.2))
ltf2 = lb2.text_frame; ltf2.word_wrap = True
for e in [
    "Product → Variation (размер, цвет, цена, остаток)",
    "Sale → SaleItem / Order → OrderItem → Variation",
    "Client → LoyaltyTransaction / Return → остаток",
    "User (ADMIN / SELLER / STOREKEEPER), Expense",
]:
    p = ltf2.paragraphs[0] if not ltf2.paragraphs[0].runs else ltf2.add_paragraph()
    p.space_before = Pt(7)
    ra = p.add_run(); ra.text = "▸  "; ra.font.size = Pt(13)
    ra.font.bold = True; ra.font.color.rgb = LBLUE
    rb = p.add_run(); rb.text = e
    rb.font.size = Pt(13); rb.font.color.rgb = DARK

rect(s9, Inches(6.25), cy + Inches(0.05), Inches(0.04), Inches(5.2), BLUE)

# Right: Results
rw = Inches(6.7); rx2 = Inches(6.45)
add_text(s9, rx2, cy + Inches(0.05), rw, Inches(0.38),
         "Достигнутые результаты", 15, bold=True, color=NAVY)

rb2 = s9.shapes.add_textbox(rx2, cy + Inches(0.48), rw, Inches(5.0))
rtf = rb2.text_frame; rtf.word_wrap = True
for title, desc in [
    ("7 модулей ERP",         "полный охват бизнес-процессов магазина"),
    ("Онлайн-витрина",        "каталог, корзина, заказы, лояльность"),
    ("28 авто-тестов",        "Jest — критическая бизнес-логика"),
    ("Финансы + Excel",       "P&L, расходы, фильтр по периоду"),
    ("POS с чеком",           "штрихкод, скидки, Z-отчёт"),
    ("Защита данных",         "JWT, rate limiting, серверная валидация цен"),
    ("Деплой PM2",            "production build, автозапуск"),
]:
    p = rtf.paragraphs[0] if not rtf.paragraphs[0].runs else rtf.add_paragraph()
    p.space_before = Pt(8)
    ra = p.add_run(); ra.text = "✓  "; ra.font.size = Pt(14)
    ra.font.bold = True; ra.font.color.rgb = LBLUE
    rb = p.add_run(); rb.text = f"{title} — "; rb.font.size = Pt(14)
    rb.font.bold = True; rb.font.color.rgb = DARK
    rc = p.add_run(); rc.text = desc
    rc.font.size = Pt(14); rc.font.color.rgb = GRAY

rect(s9, Inches(0.3), H - Inches(0.62), W - Inches(0.55), Inches(0.38), NAVY)
add_text(s9, Inches(0.3), H - Inches(0.62), W - Inches(0.55), Inches(0.38),
         "Цель достигнута: система готова к реальной эксплуатации в магазине одежды",
         15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# СЛАЙД 10 — Экономическое обоснование
# ══════════════════════════════════════════════
s10 = prs.slides.add_slide(prs.slide_layouts[6])
cy = header(s10, "Экономическое обоснование")

GREEN = RGBColor(0x16, 0xA3, 0x4A)

# ─── Левая колонка: сравнительная таблица ───
lw = Inches(6.1)
rect(s10, Inches(0.3), cy, lw, Inches(0.4), NAVY)
add_text(s10, Inches(0.4), cy + Inches(0.05), lw, Inches(0.35),
         "Сравнение: ручной учёт vs система", 13, bold=True, color=WHITE)

rows = [
    ("Показатель",              "Ручной учёт",  "Vanguard ERP"),
    ("Ввод прихода товаров",    "45–60 мин",    "5 мин"),
    ("Подготовка отчёта",       "2–3 часа",     "Автоматически"),
    ("Оформление продажи",      "3 мин",        "30 сек"),
    ("Контроль остатков",       "Раз в неделю", "В реальном времени"),
    ("Поиск клиента/истории",   "Нет",          "Мгновенно (CRM)"),
]

table_y = cy + Inches(0.45)
for i, (col1, col2, col3) in enumerate(rows):
    bg = NAVY if i == 0 else (ICE if i % 2 == 0 else WHITE)
    fc = WHITE if i == 0 else DARK
    rect(s10, Inches(0.3), table_y + i * Inches(0.57), lw, Inches(0.57), bg)
    if i == 0:
        add_text(s10, Inches(0.35), table_y + i * Inches(0.57) + Inches(0.12),
                 Inches(2.6), Inches(0.38), col1, 12, bold=True, color=WHITE)
        add_text(s10, Inches(2.95), table_y + i * Inches(0.57) + Inches(0.12),
                 Inches(1.7), Inches(0.38), col2, 12, bold=True, color=PALE, align=PP_ALIGN.CENTER)
        add_text(s10, Inches(4.65), table_y + i * Inches(0.57) + Inches(0.12),
                 Inches(1.75), Inches(0.38), col3, 12, bold=True, color=LBLUE, align=PP_ALIGN.CENTER)
    else:
        add_text(s10, Inches(0.35), table_y + i * Inches(0.57) + Inches(0.12),
                 Inches(2.6), Inches(0.38), col1, 12, color=DARK)
        add_text(s10, Inches(2.95), table_y + i * Inches(0.57) + Inches(0.12),
                 Inches(1.7), Inches(0.38), col2, 12, color=GRAY, align=PP_ALIGN.CENTER)
        rect(s10, Inches(4.6), table_y + i * Inches(0.57) + Inches(0.1),
             Inches(1.75), Inches(0.38), RGBColor(0xDC, 0xFC, 0xE7))
        add_text(s10, Inches(4.65), table_y + i * Inches(0.57) + Inches(0.12),
                 Inches(1.7), Inches(0.38), col3, 12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ─── Разделитель ───
rect(s10, Inches(6.6), cy, Inches(0.04), Inches(5.2), BLUE)

# ─── Правая колонка: расчёт экономического эффекта ───
rw = Inches(6.4); rx = Inches(6.75)
rect(s10, rx, cy, rw, Inches(0.4), BLUE)
add_text(s10, rx + Inches(0.1), cy + Inches(0.05), rw, Inches(0.35),
         "Расчёт годового эффекта", 13, bold=True, color=WHITE)

calc_box = s10.shapes.add_textbox(rx, cy + Inches(0.48), rw, Inches(2.5))
ctf = calc_box.text_frame; ctf.word_wrap = True
for label, val in [
    ("Экономия времени:", "≈ 2,5 ч/день × 250 дней = 625 ч/год"),
    ("Стоимость часа:", "400 ₽  (средняя ставка сотрудника)"),
    ("Годовая экономия:", "625 × 400 = 250 000 ₽"),
]:
    p = ctf.paragraphs[0] if not ctf.paragraphs[0].runs else ctf.add_paragraph()
    p.space_before = Pt(10)
    ra = p.add_run(); ra.text = f"{label}  "; ra.font.size = Pt(13)
    ra.font.bold = True; ra.font.color.rgb = NAVY
    rb = p.add_run(); rb.text = val
    rb.font.size = Pt(13); rb.font.color.rgb = DARK

# Итоговое число
rect(s10, rx, cy + Inches(2.4), rw, Inches(0.65), NAVY)
add_text(s10, rx + Inches(0.1), cy + Inches(2.5), rw - Inches(0.2), Inches(0.5),
         "Годовая экономия:  250 000 ₽", 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Сравнение с коммерческими аналогами
rect(s10, rx, cy + Inches(3.2), rw, Inches(0.4), NAVY)
add_text(s10, rx + Inches(0.1), cy + Inches(3.25), rw, Inches(0.35),
         "Сравнение с коммерческими аналогами", 13, bold=True, color=WHITE)

alt_box = s10.shapes.add_textbox(rx, cy + Inches(3.66), rw, Inches(1.4))
atf = alt_box.text_frame; atf.word_wrap = True
for sys, cost in [
    ("МойСклад (облако):", "от 10 800 ₽/год  (абонентская плата)"),
    ("1С: Розница:",       "от 25 000 ₽  +  ~36 000 ₽/год обслуживание"),
    ("Vanguard ERP:",      "разработка на заказ — без абонентских платежей"),
]:
    p = atf.paragraphs[0] if not atf.paragraphs[0].runs else atf.add_paragraph()
    p.space_before = Pt(8)
    ra = p.add_run(); ra.text = f"{sys}  "; ra.font.size = Pt(12)
    ra.font.bold = True; ra.font.color.rgb = LBLUE
    rb = p.add_run(); rb.text = cost
    rb.font.size = Pt(12); rb.font.color.rgb = DARK

rect(s10, Inches(0.3), H - Inches(0.62), W - Inches(0.55), Inches(0.38), BLUE)
add_text(s10, Inches(0.3), H - Inches(0.62), W - Inches(0.55), Inches(0.38),
         "Внедрение системы окупается в первый год эксплуатации",
         14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# СЛАЙД 11 — Заключение + Источники
# ══════════════════════════════════════════════
s11 = prs.slides.add_slide(prs.slide_layouts[6])
rect(s11, 0, 0, W, H, NAVY)
rect(s11, 0, 0, W, Inches(0.07), BLUE)
rect(s11, 0, H - Inches(0.07), W, Inches(0.07), BLUE)

# Заголовок
add_text(s11, 0, Inches(0.2), W, Inches(0.6),
         "Заключение", 32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s11, Inches(2.5), Inches(0.85), W - Inches(5.0), Inches(0.05), BLUE)

# Итоги — две колонки
lc = s11.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(6.2), Inches(2.8))
ltf = lc.text_frame; ltf.word_wrap = True
for item in [
    "Разработана АИС учёта товаров в магазине одежды",
    "Реализованы 7 функциональных модулей ERP",
    "28 авто-тестов + ручные сценарии тестирования",
    "Онлайн-витрина обеспечивает дополнительный канал продаж",
]:
    p = ltf.paragraphs[0] if not ltf.paragraphs[0].runs else ltf.add_paragraph()
    p.space_before = Pt(10); p.alignment = PP_ALIGN.LEFT
    ra = p.add_run(); ra.text = "▸  "; ra.font.size = Pt(15)
    ra.font.bold = True; ra.font.color.rgb = LBLUE
    rb = p.add_run(); rb.text = item
    rb.font.size = Pt(15); rb.font.color.rgb = WHITE

# Разделитель
rect(s11, Inches(6.85), Inches(0.95), Inches(0.04), Inches(2.9), BLUE)

# Перспективы — правая колонка
rc = s11.shapes.add_textbox(Inches(7.0), Inches(1.0), Inches(5.9), Inches(2.8))
rtf2 = rc.text_frame; rtf2.word_wrap = True
add_para(rtf2, "Перспективы развития:", 13, bold=True, color=PALE)
for p_item in [
    "Миграция на PostgreSQL (1 строка в Prisma schema)",
    "Мобильное приложение (React Native)",
    "Интеграция с эквайрингом",
    "Модуль инвентаризации",
]:
    p = rtf2.add_paragraph(); p.space_before = Pt(7)
    ra = p.add_run(); ra.text = "→  "; ra.font.size = Pt(13)
    ra.font.bold = True; ra.font.color.rgb = LBLUE
    rb = p.add_run(); rb.text = p_item
    rb.font.size = Pt(13); rb.font.color.rgb = WHITE

# Источники
rect(s11, Inches(0.4), Inches(3.95), W - Inches(0.8), Inches(0.04), BLUE)
add_text(s11, Inches(0.4), Inches(4.05), W - Inches(0.8), Inches(0.35),
         "Основные источники", 12, bold=True, color=PALE)

src_box = s11.shapes.add_textbox(Inches(0.4), Inches(4.4), W - Inches(0.8), Inches(1.6))
stf = src_box.text_frame; stf.word_wrap = True
for i, src in enumerate([
    "Документация Prisma ORM 5.x. — prisma.io/docs",
    "React 19 Official Docs. — react.dev",
    "Express.js Guide. — expressjs.com/guide",
    "ГОСТ Р ИСО/МЭК 12207-2010. Процессы ЖЦ ПО",
    "Фаулер М. Архитектура корпоративных программных приложений — М.: Вильямс, 2007",
    "Дейт К. Дж. Введение в системы баз данных — М.: Вильямс, 2005",
], 1):
    p = stf.paragraphs[0] if not stf.paragraphs[0].runs else stf.add_paragraph()
    p.space_before = Pt(4)
    ra = p.add_run(); ra.text = f"{i}. "; ra.font.size = Pt(11)
    ra.font.bold = True; ra.font.color.rgb = LBLUE
    rb = p.add_run(); rb.text = src
    rb.font.size = Pt(11); rb.font.color.rgb = PALE

# Автор + спасибо
add_text(s11, 0, Inches(6.2), W, Inches(0.45),
         "Созанов Георгий Валерьевич  ·  Группа ИС-22-2",
         18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s11, 0, Inches(6.7), W, Inches(0.38),
         "Спасибо за внимание! Готов ответить на вопросы.",
         14, color=LBLUE, align=PP_ALIGN.CENTER, italic=True)

# ── Save ──────────────────────────────────────
prs.save(OUT)
print(f"Готово: {OUT}")
